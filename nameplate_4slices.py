

import collections
import collections.abc
from pptx import Presentation
import pandas

df = pandas.read_excel(r'.\명단_v1.xlsx', sheet_name='Sheet1')

prs = Presentation(r'.\명찰_양식.pptx')

j = 0

# print(len(df)//4)

for i in range(0, (len(df)//4)+1) :
    add_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(add_slide_layout)

    i = i+j

    if i < len(df) :
        slide.placeholders[10].text = df.iloc[i, 2]
        slide.placeholders[11].text = df.iloc[i, 4]
        slide.placeholders[12].text = df.iloc[i, 3]
        slide.placeholders[13].text = df.iloc[i, 1]
    else :
        break

    if i+1 < len(df) :
        slide.placeholders[14].text = df.iloc[i+1, 2]
        slide.placeholders[15].text = df.iloc[i+1, 4]
        slide.placeholders[16].text = df.iloc[i+1, 3]
        slide.placeholders[17].text = df.iloc[i+1, 1]
    else :
        break 

    if i+2 < len(df) :    
        slide.placeholders[18].text = df.iloc[i+2, 2]
        slide.placeholders[19].text = df.iloc[i+2, 4]
        slide.placeholders[20].text = df.iloc[i+2, 3]
        slide.placeholders[21].text = df.iloc[i+2, 1]
    else :
        break 

    if i+3 < len(df) :
        slide.placeholders[22].text = df.iloc[i+3, 2]
        slide.placeholders[23].text = df.iloc[i+3, 4]
        slide.placeholders[24].text = df.iloc[i+3, 3]
        slide.placeholders[25].text = df.iloc[i+3, 1]
    else :
        break 

    j = j+3



# add_slide_layout = prs.slide_layouts[0]
# slide = prs.slides.add_slide(add_slide_layout)
# shapes = slide.shapes

# for shape in shapes :
#     print (str(shape.placeholder_format.idx) + " : " + shape.name)

prs.save(r'.\명찰_result.pptx')
